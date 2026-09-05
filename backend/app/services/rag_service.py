
import os
import re
import logging
import pymupdf as fitz  # PyMuPDF (new import name; `fitz` alias is deprecated)
import docx
from pptx import Presentation
os.environ.setdefault("ANONYMIZED_TELEMETRY", "False")
# ChromaDB's bundled PostHog telemetry has a known bug where its capture()
# call fails even with telemetry disabled. The failure is caught internally
# and harmless, but it logs an ERROR line on every request. Silence it here.
logging.getLogger("chromadb.telemetry.product.posthog").setLevel(logging.CRITICAL)
import chromadb
from chromadb.config import Settings as ChromaSettings
from app.config import settings
from app.services.embeddings import HashingEmbeddingFunction

_client = chromadb.PersistentClient(
    path=settings.CHROMA_DIR,
    settings=ChromaSettings(anonymized_telemetry=False),
)
_embedding_fn = HashingEmbeddingFunction()


def _collection_for(student_id: str):
    name = f"student_{student_id}".replace("-", "_")
    return _client.get_or_create_collection(name=name, embedding_function=_embedding_fn)


def extract_text(path: str, filetype: str) -> str:
    filetype = filetype.lower()
    if filetype == "pdf":
        text = []
        with fitz.open(path) as doc:
            for page in doc:
                text.append(page.get_text())
        return "\n".join(text)
    elif filetype in ("docx", "doc"):
        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    elif filetype in ("pptx", "ppt"):
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
        return "\n".join(chunks)
    elif filetype == "txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported file type: {filetype}")


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)
    return text.strip()


def chunk_text(text: str, size: int = None, overlap: int = None) -> list[str]:
    size = size or settings.CHUNK_SIZE
    overlap = overlap or settings.CHUNK_OVERLAP
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start += size - overlap
    return chunks


def ingest_document(student_id: str, document_id: str, filename: str, path: str, filetype: str) -> int:
    """Extract, clean, chunk, and store embeddings for a document. Returns chunk count."""
    raw = extract_text(path, filetype)
    cleaned = clean_text(raw)
    chunks = chunk_text(cleaned)

    if not chunks:
        return 0

    collection = _collection_for(student_id)
    ids = [f"{document_id}_{i}" for i in range(len(chunks))]
    metadatas = [{"document_id": document_id, "filename": filename, "chunk_index": i}
                 for i in range(len(chunks))]

    # Batch insert (Chroma computes embeddings automatically via its default EF)
    BATCH = 100
    for i in range(0, len(chunks), BATCH):
        collection.add(
            ids=ids[i:i + BATCH],
            documents=chunks[i:i + BATCH],
            metadatas=metadatas[i:i + BATCH],
        )
    return len(chunks)


def retrieve_context(student_id: str, query: str, document_id: str | None = None, top_k: int = None) -> list[dict]:
    top_k = top_k or settings.RETRIEVAL_TOP_K
    collection = _collection_for(student_id)
    where = {"document_id": document_id} if document_id else None
    try:
        results = collection.query(query_texts=[query], n_results=top_k, where=where)
    except Exception:
        return []

    out = []
    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    dists = results.get("distances", [[]])[0] if results.get("distances") else [None] * len(docs)
    for doc, meta, dist in zip(docs, metas, dists):
        out.append({"text": doc, "metadata": meta, "distance": dist})
    return out


def has_documents(student_id: str) -> bool:
    collection = _collection_for(student_id)
    return collection.count() > 0